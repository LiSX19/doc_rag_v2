"""
PowerPoint文档加载器

支持.pptx和.ppt格式。
添加了重试机制处理COM调用失败问题。
"""

import platform
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from .base import BaseLoader


class PPTLoader(BaseLoader):
    """PowerPoint加载器"""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        初始化PPT加载器

        Args:
            config: 配置字典，包含：
                - loader.ppt.max_retries: pywin32最大重试次数（默认3）
                - loader.ppt.retry_delay: 重试间隔（秒，默认2）
        """
        super().__init__(config)

        self._unstructured_available = self._check_unstructured()
        self._python_pptx_available = self._check_python_pptx()
        self._pywin32_available = self._check_pywin32()

        # 重试配置
        ppt_config = self.config.get('loader', {}).get('ppt', {})
        self.max_retries = ppt_config.get('max_retries', 3)
        self.retry_delay = ppt_config.get('retry_delay', 2)

    def _check_unstructured(self) -> bool:
        """检查Unstructured是否可用"""
        try:
            from unstructured.partition.pptx import partition_pptx
            from unstructured.partition.ppt import partition_ppt
            return True
        except ImportError:
            return False

    def _check_python_pptx(self) -> bool:
        """检查python-pptx是否可用"""
        try:
            from pptx import Presentation
            return True
        except ImportError:
            return False

    def _check_pywin32(self) -> bool:
        """检查pywin32是否可用（仅Windows）"""
        if platform.system() != 'Windows':
            return False
        try:
            import win32com.client
            return True
        except ImportError:
            return False

    def supports(self, file_path: Union[str, Path]) -> bool:
        """检查是否支持该文件"""
        ext = Path(file_path).suffix.lower()
        return ext in ['.pptx', '.ppt', '.ppsx']

    def load(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        加载PowerPoint文档

        Args:
            file_path: PPT文件路径

        Returns:
            包含文档内容的字典
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = file_path.suffix.lower()

        # .pptx文件
        if ext == '.pptx':
            if self._unstructured_available:
                try:
                    return self._load_with_unstructured_pptx(file_path)
                except Exception as e:
                    pass

            if self._python_pptx_available:
                return self._load_with_python_pptx(file_path)

        # .ppt文件
        elif ext in ['.ppt', '.ppsx']:
            if self._unstructured_available:
                try:
                    return self._load_with_unstructured_ppt(file_path)
                except Exception as e:
                    pass

            if self._pywin32_available:
                return self._load_with_pywin32_retry(file_path)

        raise RuntimeError(f"无法解析PowerPoint文件: {file_path}")

    def _load_with_unstructured_pptx(self, file_path: Path) -> Dict[str, Any]:
        """使用Unstructured加载.pptx"""
        from unstructured.partition.pptx import partition_pptx

        elements = partition_pptx(filename=str(file_path))

        texts = []
        slides = []
        current_slide = 1
        slide_texts = []

        for element in elements:
            text = str(element)
            if text.strip():
                texts.append(text)

                # 尝试获取幻灯片编号
                slide_number = getattr(element.metadata, 'page_number', current_slide)
                if slide_number != current_slide:
                    if slide_texts:
                        slides.append({
                            'slide_num': current_slide,
                            'content': '\n'.join(slide_texts)
                        })
                    slide_texts = []
                    current_slide = slide_number

                slide_texts.append(text)

        # 添加最后一页
        if slide_texts:
            slides.append({
                'slide_num': current_slide,
                'content': '\n'.join(slide_texts)
            })

        content = '\n\n'.join(texts)

        metadata = self.extract_metadata(file_path)
        metadata['parser'] = 'unstructured_pptx'
        metadata['slide_count'] = len(slides)

        return {
            'content': content,
            'metadata': metadata,
            'pages': slides,
        }

    def _load_with_unstructured_ppt(self, file_path: Path) -> Dict[str, Any]:
        """使用Unstructured加载.ppt"""
        from unstructured.partition.ppt import partition_ppt

        elements = partition_ppt(filename=str(file_path))

        texts = []
        for element in elements:
            text = str(element)
            if text.strip():
                texts.append(text)

        content = '\n\n'.join(texts)

        metadata = self.extract_metadata(file_path)
        metadata['parser'] = 'unstructured_ppt'

        return {
            'content': content,
            'metadata': metadata,
            'pages': [],
        }

    def _load_with_python_pptx(self, file_path: Path) -> Dict[str, Any]:
        """使用python-pptx加载.pptx"""
        from pptx import Presentation

        prs = Presentation(str(file_path))

        texts = []
        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                slide_content = '\n'.join(slide_texts)
                texts.append(slide_content)
                slides.append({
                    'slide_num': slide_num,
                    'content': slide_content
                })

        content = '\n\n'.join(texts)

        metadata = self.extract_metadata(file_path)
        metadata['parser'] = 'python_pptx'
        metadata['slide_count'] = len(slides)

        return {
            'content': content,
            'metadata': metadata,
            'pages': slides,
        }

    def _load_with_pywin32_retry(self, file_path: Path) -> Dict[str, Any]:
        """
        使用pywin32加载.ppt（带重试机制）

        Args:
            file_path: .ppt文件路径

        Returns:
            包含文档内容的字典
        """
        last_error = None

        for attempt in range(self.max_retries):
            try:
                return self._load_with_pywin32(file_path)
            except Exception as e:
                last_error = e
                error_msg = str(e)

                # 检查是否是COM错误
                if '-2147352567' in error_msg or '发生意外' in error_msg or 'RPC' in error_msg:
                    if attempt < self.max_retries - 1:
                        print(f"[PPTLoader] pywin32调用失败，{self.retry_delay}秒后重试 ({attempt + 1}/{self.max_retries})...")
                        time.sleep(self.retry_delay)
                        continue

                # 其他错误直接抛出
                raise

        # 所有重试都失败了
        raise RuntimeError(f"pywin32加载失败（已重试{self.max_retries}次）: {last_error}")

    def _load_with_pywin32(self, file_path: Path) -> Dict[str, Any]:
        """使用pywin32加载.ppt（Windows only）"""
        import win32com.client

        powerpoint = None
        presentation = None

        try:
            # 创建PowerPoint应用实例
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")

            # 尝试设置可见性（某些情况下可能失败）
            try:
                powerpoint.Visible = False
            except:
                pass  # 忽略设置可见性失败

            # 禁用警告和自动更正
            try:
                powerpoint.DisplayAlerts = False
            except:
                pass

            # 打开演示文稿
            presentation = powerpoint.Presentations.Open(
                str(file_path.absolute()),
                ReadOnly=True,
                Untitled=False,
                WithWindow=False
            )

            texts = []
            slides = []

            for slide_num in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_num)

                slide_texts = []
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        text_frame = shape.TextFrame
                        if text_frame.HasText:
                            slide_texts.append(text_frame.TextRange.Text)

                if slide_texts:
                    slide_content = '\n'.join(slide_texts)
                    texts.append(slide_content)
                    slides.append({
                        'slide_num': slide_num,
                        'content': slide_content
                    })

            content = '\n\n'.join(texts)

            metadata = self.extract_metadata(file_path)
            metadata['parser'] = 'pywin32'
            metadata['slide_count'] = len(slides)

            return {
                'content': content,
                'metadata': metadata,
                'pages': slides,
            }

        finally:
            # 确保关闭演示文稿和应用
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass

            if powerpoint:
                try:
                    powerpoint.Quit()
                except:
                    pass

            # 强制垃圾回收，释放COM对象
            import gc
            gc.collect()
