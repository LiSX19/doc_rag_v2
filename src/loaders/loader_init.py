#!/usr/bin/env python3
"""
Unstructured模型预加载脚本

在处理文档前运行此脚本，预先加载unstructured所需的模型，
避免首次解析文档时的长时间等待。

使用方法:
    python -m src.loaders.loader_init
"""

import os
import sys

os.environ.setdefault('UNSTRUCTURED_DISABLE_VISION', '0')


def preload_unstructured():
    """预加载unstructured核心模块和模型"""
    print("=" * 50)
    print("Unstructured 模型预加载")
    print("=" * 50)

    print("\n[1/5] 检查核心模块...")
    try:
        import unstructured
        print(f"  ✓ unstructured {unstructured.__version__}")
    except ImportError as e:
        print(f"  ✗ unstructured导入失败: {e}")
        return False

    print("\n[2/5] 加载PDF解析模块...")
    try:
        from unstructured.partition.pdf import partition_pdf
        print("  ✓ partition_pdf 已加载")
    except ImportError as e:
        print(f"  ✗ partition_pdf导入失败: {e}")

    print("\n[3/5] 加载Word(.docx)解析模块...")
    try:
        from unstructured.partition.docx import partition_docx
        print("  ✓ partition_docx 已加载")
    except ImportError as e:
        print(f"  ✗ partition_docx导入失败: {e}")

    print("\n[4/5] 加载Excel(.xlsx)解析模块...")
    try:
        from unstructured.partition.xlsx import partition_xlsx
        print("  ✓ partition_xlsx 已加载")
    except ImportError as e:
        print(f"  ✗ partition_xlsx导入失败: {e}")

    print("\n[5/5] 加载PPT解析模块...")
    try:
        from unstructured.partition.pptx import partition_pptx
        print("  ✓ partition_pptx 已加载")
    except ImportError as e:
        print(f"  ✗ partition_pptx导入失败: {e}")

    print("\n" + "=" * 50)
    print("预加载完成！")
    print("=" * 50)
    print("""
提示: 首次运行可能会下载模型，耗时较长。
      之后再次运行将使用缓存，速度会显著提升。

如果模型下载过慢，可以设置镜像源:
    export HF_ENDPOINT=https://hf-mirror.com
""")
    return True


def warm_up_models():
    """使用虚拟输入触发模型实际加载"""
    print("\n正在进行模型热身...")

    import tempfile
    from pathlib import Path

    # 创建虚拟文本文件
    dummy_text = Path(tempfile.mktemp(suffix='.txt'))
    dummy_text.write_text("warm up text", encoding='utf-8')

    # 创建虚拟docx文件
    dummy_docx = Path(tempfile.mktemp(suffix='.docx'))
    try:
        from docx import Document
        doc = Document()
        doc.add_paragraph("warm up docx")
        doc.save(str(dummy_docx))
    except ImportError:
        print("  ⚠ python-docx未安装，跳过docx热身")

    # 创建虚拟xlsx文件
    dummy_xlsx = Path(tempfile.mktemp(suffix='.xlsx'))
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = 'warm up xlsx'
        wb.save(str(dummy_xlsx))
    except ImportError:
        print("  ⚠ openpyxl未安装，跳过xlsx热身")

    # 1. 热身text模块
    print("  [1/4] 热身文本解析...")
    try:
        from unstructured.partition.text import partition_text
        elements = list(partition_text(filename=str(dummy_text)))
        print(f"  ✓ 文本解析成功 ({len(elements)}个元素)")
    except Exception as e:
        print(f"  ⚠ 文本解析热身失败: {e}")

    # 2. 热身docx模块
    print("  [2/4] 热身Word解析...")
    if dummy_docx.exists():
        try:
            from unstructured.partition.docx import partition_docx
            elements = list(partition_docx(filename=str(dummy_docx), languages=['chi_sim']))
            print(f"  ✓ Word解析成功 ({len(elements)}个元素)")
        except Exception as e:
            print(f"  ⚠ Word解析热身失败: {e}")

    # 3. 热身xlsx模块
    print("  [3/4] 热身Excel解析...")
    if dummy_xlsx.exists():
        try:
            from unstructured.partition.xlsx import partition_xlsx
            elements = list(partition_xlsx(filename=str(dummy_xlsx)))
            print(f"  ✓ Excel解析成功 ({len(elements)}个元素)")
        except Exception as e:
            print(f"  ⚠ Excel解析热身失败: {e}")

    # 4. 热身pptx模块
    print("  [4/4] 热身PowerPoint解析...")
    try:
        from unstructured.partition.pptx import partition_pptx
        # 创建简单的pptx文件
        from pptx import Presentation
        dummy_pptx = Path(tempfile.mktemp(suffix='.pptx'))
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Warm Up"
        prs.save(str(dummy_pptx))
        
        elements = list(partition_pptx(filename=str(dummy_pptx), languages=['chi_sim']))
        print(f"  ✓ PowerPoint解析成功 ({len(elements)}个元素)")
        dummy_pptx.unlink()
    except Exception as e:
        print(f"  ⚠ PowerPoint解析热身失败: {e}")

    # 清理临时文件
    for f in [dummy_text, dummy_docx, dummy_xlsx]:
        if f.exists():
            try:
                f.unlink()
            except:
                pass

    print("\n模型热身完成！现在可以开始文档处理。")
    print("\n提示: 模型已缓存到本地，后续运行将显著加快。")


def main():
    """主函数"""
    preload_unstructured()

    import argparse
    parser = argparse.ArgumentParser(description='Unstructured模型预加载')
    parser.add_argument('--no-warm-up', action='store_true',
                        help='跳过模型热身（不推荐）')
    args = parser.parse_args()

    if not args.no_warm_up:
        warm_up_models()


if __name__ == '__main__':
    main()
